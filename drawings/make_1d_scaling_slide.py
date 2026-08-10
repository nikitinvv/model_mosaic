"""Single-slide PPTX explaining how tomo + prop scale to whole-brain volumes
by reducing every 2-D operator into chains of 1-D computations.

Run: python make_1d_scaling_slide.py  →  drawings/1d_scaling.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# --- palette ---------------------------------------------------------------
BG          = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_C     = RGBColor(0x11, 0x24, 0x3D)
ACCENT_C    = RGBColor(0x1F, 0x77, 0xB4)   # blue
ACCENT2_C   = RGBColor(0xE1, 0x5F, 0x2E)   # orange
BODY_C      = RGBColor(0x33, 0x33, 0x33)
GRID_C      = RGBColor(0xEE, 0xEE, 0xEE)
STRIKE_C    = RGBColor(0xD3, 0x28, 0x2F)   # red for "gone"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=BODY_C, align=PP_ALIGN.LEFT,
             mono=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if mono:
        r.font.name = 'Consolas'
    return tb


def add_rich_line(slide, x, y, w, h, runs, size=14, mono=False, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP):
    """runs = list of (text, {size?, bold?, color?, mono?})"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.color.rgb = opts.get('color', BODY_C)
        r.font.name = 'Consolas' if opts.get('mono', mono) else 'Calibri'
    return tb


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_pass_chain(slide, x, y, w, h, boxes, arrow_c=ACCENT_C,
                   box_fill=RGBColor(0xF2, 0xF6, 0xFC),
                   box_line=ACCENT_C, mono=True):
    """Draw a row of small labeled boxes with arrows between."""
    n = len(boxes)
    gap = Inches(0.15)
    total_gap = gap * (n - 1)
    each = (w - total_gap) // n
    xi = x
    for i, label in enumerate(boxes):
        sh = add_box(slide, xi, y, each, h, fill=box_fill,
                     line=box_line, line_w=0.75, radius=0.15)
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(11)
        r.font.name = 'Consolas' if mono else 'Calibri'
        r.font.color.rgb = TITLE_C
        r.font.bold = True
        if i < n - 1:
            ax = xi + each
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           ax, y + h // 2 - Inches(0.10), gap, Inches(0.20))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = arrow_c
            arrow.line.fill.background()
            arrow.shadow.inherit = False
        xi += each + gap


# --- build the slide -------------------------------------------------------
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# Title
add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
         'Scaling tomo + prop to whole-brain volumes',
         size=32, bold=True, color=TITLE_C)
add_text(slide, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.40),
         'Key idea:  every 2-D operator is separable → we only run 1-D computations',
         size=17, bold=False, color=ACCENT2_C)


# ---------------- Tomo block ----------------
tomo_y0 = Inches(1.55)
add_text(slide, Inches(0.5), tomo_y0, Inches(6.0), Inches(0.35),
         'Tomo — USFFT Radon', size=18, bold=True, color=ACCENT_C)

add_pass_chain(slide, Inches(0.5), tomo_y0 + Inches(0.45),
               Inches(6.15), Inches(0.42),
               ['obj', '·φ', 'FFT_x', 'FFT_y', 'gather', 'IFFT_r', 'sino'])

tomo_notes = [
    ('φ(x, y)', '='), '  φ₁ᴅ(x) · φ₁ᴅ(y)', '     → stored as 2 vectors of length n',
    ('c2d(x, y)', '='), '  c2d₁ᴅ(x) · c2d₁ᴅ(y)', '  → 2 vectors of length 2n',
    ('(x₀, y₀)', '='), '  (cos θₖ·r, −sin θₖ·r)', '  → recomputed inside gather kernel',
]

# Render three rows
lines_tomo = [
    ("φ(x, y) = ", "φ₁ᴅ(x) · φ₁ᴅ(y)",     "  → stored as 2 vectors of length n"),
    ("c2d(x, y) = ", "c2d₁ᴅ(x) · c2d₁ᴅ(y)", "  → 2 vectors of length 2n"),
    ("(x₀, y₀) = ", "(cos θₖ · r, −sin θₖ · r)", "  → recomputed inside CUDA gather kernel"),
]
row_y = tomo_y0 + Inches(1.05)
for lbl, factored, note in lines_tomo:
    add_rich_line(slide, Inches(0.55), row_y, Inches(6.3), Inches(0.28),
                  [(lbl, {'mono': True, 'size': 12, 'color': BODY_C}),
                   (factored, {'mono': True, 'size': 12, 'bold': True, 'color': ACCENT_C}),
                   (note, {'size': 11, 'color': BODY_C})])
    row_y += Inches(0.28)


# ---------------- Prop block ----------------
prop_y0 = Inches(1.55)
add_text(slide, Inches(7.0), prop_y0, Inches(6.0), Inches(0.35),
         'Prop — Fresnel angular-spectrum', size=18, bold=True, color=ACCENT_C)

add_pass_chain(slide, Inches(7.0), prop_y0 + Inches(0.45),
               Inches(6.15), Inches(0.42),
               ['ψ', 'FFT_x', '·K_x(fx)', 'FFT_y', '·K_y(fy)', 'IFFT_y', 'IFFT_x'])

# Equation + note
add_rich_line(slide, Inches(7.05), prop_y0 + Inches(1.05), Inches(6.2), Inches(0.32),
              [('K(fx, fy) = exp(−iπλL(fx² + fy²)) = ',
                {'mono': True, 'size': 13}),
               ('K_x(fx) · K_y(fy)',
                {'mono': True, 'size': 13, 'bold': True, 'color': ACCENT_C})])
add_text(slide, Inches(7.05), prop_y0 + Inches(1.40), Inches(6.2), Inches(0.32),
         'Fresnel kernel is separable in the Fourier domain — stored as 2 vectors of length 2n.',
         size=12)


# ---------------- Eliminated (both) ----------------
elim_y = Inches(3.85)
add_text(slide, Inches(0.5), elim_y, Inches(12.3), Inches(0.35),
         'Never allocated (would dominate at large N):',
         size=15, bold=True, color=STRIKE_C)

elim_rows = [
    ('Tomo', 'φ  (n × n) c64',           'n² ·  8 B',   '2 vectors × n  c64  (~ MB)'),
    ('Tomo', 'c2dfftshift  (2n × 2n) i8', '(2n)² · 1 B', '2 vectors × 2n i8  (~ KB)'),
    ('Tomo', 'x, y  sample tables  f32', '2 · nθ · n · 4 B', 'recomputed in kernel from (cos θ, sin θ, idx)'),
    ('Prop', 'K  (2n × 2n) c64',         '(2n)² · 8 B', '2 vectors × 2n c64  (~ MB)'),
]

# header row
header_y = elim_y + Inches(0.4)
hdr = [('who', 1.0), ('2-D object', 3.4), ('naive size', 2.3), ('what we store instead', 5.1)]
xh = Inches(0.5)
for h_text, w in hdr:
    add_text(slide, xh, header_y, Inches(w), Inches(0.28),
             h_text, size=11, bold=True, color=TITLE_C)
    xh += Inches(w + 0.1)

for i, (who, obj, naive, actual) in enumerate(elim_rows):
    yy = header_y + Inches(0.30 + i * 0.28)
    # zebra stripe
    if i % 2 == 0:
        add_box(slide, Inches(0.45), yy - Inches(0.02),
                Inches(12.4), Inches(0.26),
                fill=GRID_C, line=None)
    xx = Inches(0.5)
    add_text(slide, xx, yy, Inches(1.0), Inches(0.28), who,
             size=11, bold=True, color=ACCENT2_C if who == 'Prop' else ACCENT_C)
    xx += Inches(1.1)
    add_text(slide, xx, yy, Inches(3.4), Inches(0.28), obj, size=11, mono=True)
    xx += Inches(3.5)
    add_text(slide, xx, yy, Inches(2.3), Inches(0.28), naive, size=11, mono=True,
             color=STRIKE_C)
    xx += Inches(2.4)
    add_text(slide, xx, yy, Inches(5.1), Inches(0.28), actual, size=11)


# ---------------- Second-order tricks ----------------
tricks_y = Inches(5.85)
add_text(slide, Inches(0.5), tricks_y, Inches(12.3), Inches(0.35),
         'Second-order tricks that ride the 1-D structure',
         size=15, bold=True, color=TITLE_C)

tricks = [
    ('Host-staged 2-D grid',
     'the full 2N × 2N Fourier grid lives PINNED on host; GPU sees only per-pass strips',
     'GPU peak fixed at ~35 GB regardless of N'),
    ('Banded pinned memory',
     'one 576 GB cudaHostAlloc  →  4–8 blocks of ≤ 144 GB each',
     'fits under NVIDIA driver’s per-alloc cap'),
    ('Streamed 3-way pipe',
     'load ‖ compute ‖ store per chunk with 2× ping-pong pinned buffers',
     'overlap hides H2D / D2H latency'),
]
tw = Inches(4.15)
th = Inches(1.25)
tx = Inches(0.5)
ty = tricks_y + Inches(0.4)
gap = Inches(0.12)
for name, what, payoff in tricks:
    add_box(slide, tx, ty, tw, th, fill=RGBColor(0xF7, 0xF7, 0xF7),
            line=RGBColor(0xCC, 0xCC, 0xCC), radius=0.10)
    add_text(slide, tx + Inches(0.15), ty + Inches(0.08), tw - Inches(0.3), Inches(0.28),
             name, size=13, bold=True, color=ACCENT_C)
    add_text(slide, tx + Inches(0.15), ty + Inches(0.38), tw - Inches(0.3), Inches(0.60),
             what, size=11, color=BODY_C)
    add_text(slide, tx + Inches(0.15), ty + Inches(0.98), tw - Inches(0.3), Inches(0.32),
             '→ ' + payoff, size=11, bold=True, color=ACCENT2_C)
    tx += tw + gap


# ---------------- Save --------------------------------------------------------
here = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(here, '1d_scaling.pptx')
prs.save(out)
print('saved:', out)
